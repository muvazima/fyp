import streamlit as st
import numpy as np
import pandas as pd
import os
import pickle
import re
from sklearn.datasets import load_files
import glob
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from summarizer import Summarizer
from pptx import Presentation 
from pptx.util import Inches, Pt  
import base64
import win32com.client

model = Summarizer('distilbert-base-uncased')  


def file_selector(folder_path='./Parsed_Papers'):
    filenames = os.listdir(folder_path)
    selected_filename = st.sidebar.selectbox('Select a file', filenames)
    return os.path.join(folder_path, selected_filename)

def gen_dict(uploaded_file):

	text = str(uploaded_file.read())
	text = re.sub("\n", " ", text)
	text=re.sub("\t"," ",text)
	text=re.sub("\r\n"," ",text)
	text = (text.encode('ascii', 'ignore')).decode("utf-8")
	text= re.sub(' +', ' ', text)
	sections = re.findall(r"@&#\w+@&#", text)
	main_title = re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)

	d={}
	d['@&#MAIN-TITLE@&#']=re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)
	for i in sections:
  		d[i]=re.findall(i+'(.*?)@&#', text,  flags = re.I)

	del d['@&#REFERENCES@&#'] 
	return d

def read_file(filename):

	f = open(filename, 'r', encoding="utf-8")
	text = str(f.read())
	f.close()
	text = re.sub("\n", " ", text)
	text=re.sub("\t"," ",text)
	text = (text.encode('ascii', 'ignore')).decode("utf-8")
	text= re.sub(' +', ' ', text)
	sections = re.findall(r"@&#\w+@&#", text)
	main_title = re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)

	d={}
	d['@&#MAIN-TITLE@&#']=re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)
	for i in sections:
  		d[i]=re.findall(i+'(.*?)@&#', text,  flags = re.I)

	del d['@&#REFERENCES@&#'] 
	return d

def summarize(d):

	donotsummarize=['MAIN-TITLE','HIGHLIGHTS','KEYPHRASES','REFERENCES','ACKNOWLEDGEMENTS']
	lines=[]
	dclean={}
	for i in d:
		iclean=re.sub("@&#","", i)
		if iclean in donotsummarize:
			lines.append(iclean+": "+d[i][0]+"\n\n\n")
			dclean[iclean]=d[i][0]
		else:
			st=model(d[i][0])
			lines.append(iclean+": "+st+"\n\n\n")
			dclean[iclean]=st
	return dclean

def create_ppt(dclean,filename):

	prs = Presentation() 
	first_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(first_slide_layout) 
	slide.shapes.title.text = dclean['MAIN-TITLE']

	for i in dclean: 
  		bullet_slide_layout = prs.slide_layouts[1]
  		slide = prs.slides.add_slide(bullet_slide_layout)
  		shapes = slide.shapes
  		title_shape = shapes.title
  		body_shape = shapes.placeholders[1]
  		title_shape.text = i
  		l=dclean[i].split('.')
  		tf = body_shape.text_frame
  		tf.text = l[0]

  		for j in l[1:len(l)]:
  			p = tf.add_paragraph()
  			p.text =j

	#prs.save('/Users/Manam/fyp/PPTs/'+filename+'.pptx')
	prs.save(filename)




def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

def main():

	st.sidebar.title("Paper to PPT converter")
	filename = file_selector()
	st.sidebar.write('You selected `%s`' % filename)
	d=read_file(filename)
	dclean=summarize(d)

	filename=st.sidebar.text_input("Enter File Path to save PPT")
	#uploaded_file = st.sidebar.file_uploader("Choose a file", type=['txt'])
	if filename!='':
		#d=gen_dict(uploaded_file)
		#print(d)
		#dclean=summarize(d)
		#print(dclean)
		#filename=st.sidebar.text_input("Enter File Path to save PPT")
		create_ppt(dclean,filename)
		st.header("pptx file saved.")
		#inputFileName='/Users/Manam/fyp/PPTs/'+filename+'.pptx'
		#outputFileName='/Users/Manam/fyp/PDFs/'+filename+'.pdf'
		inputFileName=filename
		outputFileName=filename

		PPTtoPDF(inputFileName, outputFileName)
		if outputFileName[-3:] != 'pdf':
                        outputFileName = outputFileName + ".pdf"
		with open(outputFileName,"rb") as f:
			base64_pdf = base64.b64encode(f.read()).decode('utf-8')
		pdf_display = f'<embed src="data:application/pdf;base64,{base64_pdf}" width="700" height="1000" type="application/pdf">'
		selection = st.sidebar.radio("Select voice", ['Muvazima','Srikanth','Amitabh Bachchan'])

		st.markdown(pdf_display, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
