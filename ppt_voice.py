import streamlit as st
from synthesizer.inference import Synthesizer
from encoder import inference as encoder
from vocoder import inference as vocoder
from pathlib import Path
import numpy as np
import soundfile as sf
import os
import librosa
import glob
from helper import draw_embed, create_spectrogram, read_audio, record, save_record
import os
import pickle
import re
from sklearn.datasets import load_files
import glob
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from summarizer import Summarizer
from pptx import Presentation 
from pptx.util import Inches, Pt  
import base64
#import win32com.client

model = Summarizer('distilbert-base-uncased')  


def file_selector(folder_path='./Parsed_Papers'):
    filenames = os.listdir(folder_path)
    selected_filename = st.sidebar.selectbox('Select a file', filenames)
    return os.path.join(folder_path, selected_filename)

def gen_dict(uploaded_file):

    text = str(uploaded_file.read())
    text = re.sub("\n", " ", text)
    text=re.sub("\t"," ",text)
    text=re.sub("\r\n"," ",text)
    text = (text.encode('ascii', 'ignore')).decode("utf-8")
    text= re.sub(' +', ' ', text)
    sections = re.findall(r"@&#\w+@&#", text)
    
    main_title = re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)

    d={}
    d['@&#MAIN-TITLE@&#']=re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)
    for i in sections:
        d[i]=re.findall(i+'(.*?)@&#', text,  flags = re.I)

    del d['@&#REFERENCES@&#'] 
    return d

def read_file(filename):

    f = open(filename, 'r', encoding="utf-8")
    text = str(f.read())
    f.close()
    text = re.sub("\n", " ", text)
    text=re.sub("\t"," ",text)
    text = (text.encode('ascii', 'ignore')).decode("utf-8")
    text= re.sub(' +', ' ', text)
    sections = re.findall(r"@&#\w+@&#", text)
    main_title = re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)

    d={}
    d['@&#MAIN-TITLE@&#']=re.findall(r'@&#MAIN-TITLE@&#(.*?)@&#', text, flags = re.I)
    for i in sections:
        d[i]=re.findall(i+'(.*?)@&#', text,  flags = re.I)

    del d['@&#REFERENCES@&#'] 
    return d

#@st.cache(allow_output_mutation=True)
def summarize(d):

    donotsummarize=['MAIN-TITLE','HIGHLIGHTS','KEYPHRASES','REFERENCES','ACKNOWLEDGEMENTS']
    #lines=[]
    dclean={}
    for i in d:
        iclean=re.sub("@&#","", i)
        if iclean in donotsummarize:
            #lines.append(iclean+": "+d[i][0]+"\n\n\n")
            dclean[iclean]=d[i][0]
        else:
            st=model(d[i][0])
            #lines.append(iclean+": "+st+"\n\n\n")
            dclean[iclean]=st
    return dclean

@st.cache(show_spinner=False)
def create_ppt_new(dclean,filename):

    prs = Presentation() 
    first_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(first_slide_layout) 
    slide.shapes.title.text = dclean['MAIN-TITLE']

    for i in dclean: 
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = i
        l=dclean[i].split('.')
        tf = body_shape.text_frame
        tf.text = l[0]
        words=len(l[0])- l[0].count(' ')
        for j in l[1:len(l)-1]:
        	if(words>20):
        		bullet_slide_layout = prs.slide_layouts[1]
        		slide = prs.slides.add_slide(bullet_slide_layout)
        		shapes = slide.shapes
        		#title_shape = shapes.title
        		body_shape = shapes.placeholders[1]
        		tf = body_shape.text_frame
        		#tf.text = j
        		p = tf.add_paragraph()
        		p.text =j
        		words=len(j)- j.count(' ')
        	else:
        		p = tf.add_paragraph()
        		p.text =j
        		words=words+(len(j)- j.count(' '))

    prs.save('/Users/Manam/final-fyp/PPTs/'+filename+'.pptx')
    #prs.save(filename)


@st.cache(show_spinner=False)
def create_ppt(dclean,filename):

    prs = Presentation() 
    first_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(first_slide_layout) 
    slide.shapes.title.text = dclean['MAIN-TITLE']

    for i in dclean[1:]: 
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = i
        l=dclean[i].split('.')
        tf = body_shape.text_frame
        tf.text = l[0]
        #words=len(l[0])
        for j in l[1:]:
            p = tf.add_paragraph()
            p.text =j

    prs.save('/Users/Manam/final-fyp/PPTs/'+filename+'.pptx')
    #prs.save(filename)



# @st.cache
# def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
#     powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
#     powerpoint.Visible = 1

#     if outputFileName[-3:] != 'pdf':
#         outputFileName = outputFileName + ".pdf"
#     deck = powerpoint.Presentations.Open(inputFileName)
#     deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
#     deck.Close()
#     powerpoint.Quit()

@st.cache(show_spinner=False)
def create_embedding(selected_filename):
    in_fpath = Path(selected_filename.replace('"', "").replace("'", ""))
    original_wav, sampling_rate = librosa.load(str(in_fpath))
    preprocessed_wav = encoder.preprocess_wav(original_wav, sampling_rate)
    embed = encoder.embed_utterance(preprocessed_wav)
    return embed

#"# Streamlit showcase"
#@st.cache(allow_output_mutation=True,hash_funcs={FooType: hash})
def load_model(MODEL_PATH):
    model = load_model(MODEL_PATH)
    model._make_predict_function()
    model.summary()  # included to make it visible when model is reloaded
    return model

#model_load_state = st.text("Loading pretrained models...")

seed = 42
low_mem = False
num_generated = 0
enc_model_fpath = Path("encoder/saved_models/pretrained.pt")
syn_model_dir = Path("synthesizer/saved_models/logs-pretrained")
voc_model_fpath = Path("vocoder/saved_models/pretrained/pretrained.pt")
encoder.load_model(enc_model_fpath)
#encoder=load_model(enc_model_fpath)
#vocoder=load_model(voc_model_fpath)
synthesizer = Synthesizer(
    syn_model_dir.joinpath("taco_pretrained"), low_mem=low_mem, seed=seed
)
#synthesizer.load_model('synthesizer/saved_models/logs-pretrained/pretrained.pt')
vocoder.load_model(voc_model_fpath)

#model_load_state.text("Loaded pretrained models!")

def main():

    st.title("Paper to PPT converter with custom voice delivery")
    filename = file_selector()
    st.sidebar.write('You selected `%s`' % filename)
    d=read_file(filename)
    dclean=summarize(d)

    filename=st.sidebar.text_input("Enter PPT file name")
    #filename='./PPTs/'
    #uploaded_file = st.sidebar.file_uploader("Choose a file", type=['txt'])
    if filename!='':
        #d=gen_dict(uploaded_file)
        #print(d)
        #dclean=summarize(d)
        #print(dclean)
        #filename=st.sidebar.text_input("Enter File Path to save PPT")
        create_ppt_new(dclean,filename)
        st.header("pptx file saved.")
        #inputFileName='/Users/Manam/fyp/PPTs/'+filename+'.pptx'
        #outputFileName='/Users/Manam/fyp/PDFs/'+filename+'.pdf'
        inputFileName=filename
        outputFileName=filename

        # PPTtoPDF(inputFileName, outputFileName)
        # if outputFileName[-3:] != 'pdf':
        #                 outputFileName = outputFileName + ".pdf"
        # with open(outputFileName,"rb") as f:
        #     base64_pdf = base64.b64encode(f.read()).decode('utf-8')
        # pdf_display = f'<embed src="data:application/pdf;base64,{base64_pdf}" width="700" height="1000" type="application/pdf">'

        # st.markdown(pdf_display, unsafe_allow_html=True)

        st.sidebar.title("Record your own voice")

        audiofilename = st.sidebar.text_input("Enter a filename for your voice: ")

        if st.sidebar.button(f"Click to Record"):
            if audiofilename == "":
                st.warning("Choose a filename.")
            else:
                record_state = st.text("Recording...")
                duration = 10  # seconds
                fs = 48000
                myrecording = record(duration, fs)
                record_state.text(f"Saving sample as {filename}.mp3")

                path_myrecording = f"./samples/{filename}.mp3"

                save_record(path_myrecording, myrecording, fs)
                #record_state.text(f"Done! Saved sample as {filename}.mp3")

                st.sidebar.audio(read_audio(path_myrecording))

                #fig = create_spectrogram(path_myrecording)
                #st.pyplot(fig)

        audio_folder = "samples"
        filenames = glob.glob(os.path.join(audio_folder, "*.mp3"))
        selected_filename = st.sidebar.selectbox("Select a voice", filenames)
        
        if selected_filename is not None:
    # Create embedding
            embed=create_embedding(selected_filename)
    #st.success("Created the embedding")

    #st.audio(read_audio(in_fpath))

    #fig = draw_embed(embed, "myembedding", None)
    #st.pyplot(fig)
        for i in dclean:
            text = dclean[i]
            texts = [text]
            embeds = [embed]

    # generate waveform
            #with st.spinner("Generating your speech..."):
            specs = synthesizer.synthesize_spectrograms(texts, embeds)
            spec = specs[0]
            
        
            generated_wav = vocoder.infer_waveform(spec)
            generated_wav = np.pad(generated_wav, (0, synthesizer.sample_rate), mode="constant")
            generated_wav = encoder.preprocess_wav(generated_wav)
    

    # Save it on the disk
            opfilename = "Output/"+filename+"-"+i+".wav"
            sf.write(opfilename, generated_wav.astype(np.float32), synthesizer.sample_rate)
            st.text(i)
            st.audio(read_audio(opfilename))
            #st.audio(read_audio(filename))

if __name__ == "__main__":
    main()